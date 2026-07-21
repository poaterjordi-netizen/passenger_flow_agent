from __future__ import annotations

from pathlib import Path
import tempfile
from typing import Any, cast
from xml.etree import ElementTree
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "assets" / "metro-passenger-flow-agent-overview-cn.pptx"
SCREENSHOT = ROOT / "docs" / "assets" / "assistant-web-ui.png"

W = 13.333
H = 7.5
FONT = "Microsoft YaHei"
NAVY = "18354A"
BLUE = "087EA4"
CYAN = "00A6C8"
PALE = "EAF6F9"
MIST = "F5F9FB"
WHITE = "FFFFFF"
TEXT = "173042"
MUTED = "607785"
GREEN = "1E8E69"
ORANGE = "D97706"
RED = "C2413A"
BORDER = "C9DCE4"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = TEXT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin: float = 0.06,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str = BORDER,
    radius: bool = True,
    line_width: float = 1,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    if radius:
        try:
            shape.adjustments[0] = 0.08
        except (IndexError, ValueError):
            pass
    return shape


def add_circle(slide, label: str, x: float, y: float, d: float, *, fill: str = CYAN):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(fill)
    circle.line.color.rgb = rgb(fill)
    add_text(slide, label, x, y, d, d, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    return circle


def add_line(slide, x1: float, y1: float, x2: float, y2: float, *, color: str = CYAN, width: float = 2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_header(slide, title: str, section: str, page: int):
    add_text(slide, section, 0.55, 0.2, 2.4, 0.3, size=10, color=CYAN, bold=True)
    add_text(slide, title, 0.55, 0.55, 12.0, 0.52, size=27, color=NAVY, bold=True)
    add_line(slide, 0.55, 1.17, 12.78, 1.17, color=CYAN, width=1.4)
    add_text(slide, str(page), 12.25, 0.18, 0.5, 0.35, size=12, color=BLUE, bold=True, align=PP_ALIGN.RIGHT)


def add_footer(slide, text: str = "本地受治理原型｜合成数据与生产验收分开"):
    add_text(slide, text, 0.55, 7.13, 8.5, 0.22, size=8.5, color=MUTED)
    add_text(slide, "Metro Passenger Flow Agent", 9.6, 7.13, 3.15, 0.22, size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, title: str, body: str, x: float, y: float, w: float, h: float, *, accent: str = CYAN):
    add_rect(slide, x, y, w, h, fill=WHITE, line=BORDER)
    add_rect(slide, x, y, 0.08, h, fill=accent, line=accent, radius=False, line_width=0)
    add_text(slide, title, x + 0.23, y + 0.18, w - 0.42, 0.42, size=16, color=NAVY, bold=True)
    add_text(slide, body, x + 0.23, y + 0.67, w - 0.42, h - 0.85, size=12.5, color=TEXT, valign=MSO_ANCHOR.TOP)


def add_pill(slide, text: str, x: float, y: float, w: float, *, fill: str = PALE, color: str = BLUE):
    add_rect(slide, x, y, w, 0.38, fill=fill, line=fill)
    add_text(slide, text, x, y, w, 0.38, size=10.5, color=color, bold=True, align=PP_ALIGN.CENTER)


def style_table(table) -> None:
    for row_index, row in enumerate(table.rows):
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(NAVY if row_index == 0 else WHITE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = FONT
                    run.font.size = Pt(11 if row_index == 0 else 10.5)
                    run.font.bold = row_index == 0
                    run.font.color.rgb = rgb(WHITE if row_index == 0 else TEXT)


def normalize_chart_axis_ids(path: Path) -> None:
    """Normalize python-pptx chart axis IDs to a portable positive Int32 range."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as handle:
        temporary = Path(handle.name)
    try:
        with ZipFile(path, "r") as source, ZipFile(temporary, "w", ZIP_DEFLATED) as target:
            for info in source.infolist():
                payload = source.read(info.filename)
                if info.filename.startswith("ppt/charts/chart") and info.filename.endswith(".xml"):
                    # The XML is generated in this process, not supplied by an untrusted source.
                    root = ElementTree.fromstring(payload)
                    changed = False
                    for element in root.iter():
                        if element.tag.rsplit("}", 1)[-1] not in {"axId", "crossAx"}:
                            continue
                        value = element.get("val")
                        if value is not None:
                            normalized = int(value) % (2**31 - 1)
                            normalized = normalized or 1
                            element.set("val", str(normalized))
                            changed = True
                    if changed:
                        payload = ElementTree.tostring(root, encoding="utf-8", xml_declaration=True)
                elif (
                    info.filename.startswith("ppt/slideMasters/slideMaster")
                    or info.filename.startswith("ppt/slideLayouts/slideLayout")
                ) and info.filename.endswith(".xml"):
                    # Safe here because this OOXML package was generated above in this process.
                    # Remove blank-template date/footer/slide-number placeholders. They otherwise
                    # appear as stale defaults such as "1/27/13" or duplicate page numbers.
                    root = ElementTree.fromstring(payload)
                    changed = False
                    for parent in root.iter():
                        for child in list(parent):
                            if child.tag.rsplit("}", 1)[-1] != "sp":
                                continue
                            placeholders = [
                                node
                                for node in child.iter()
                                if node.tag.rsplit("}", 1)[-1] == "ph"
                            ]
                            if any(
                                node.get("type") in {"dt", "ftr", "sldNum"}
                                for node in placeholders
                            ):
                                parent.remove(child)
                                changed = True
                    if changed:
                        payload = ElementTree.tostring(root, encoding="utf-8", xml_declaration=True)
                target.writestr(info, payload)
        temporary.replace(path)
    finally:
        temporary.unlink(missing_ok=True)


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # 1 封面
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)
    add_rect(slide, 0, 0, 4.35, H, fill=NAVY, line=NAVY, radius=False, line_width=0)
    add_text(slide, "受治理的\n地铁客流智能体", 0.62, 1.15, 3.2, 1.65, size=31, color=WHITE, bold=True, valign=MSO_ANCHOR.TOP)
    add_text(slide, "让大模型负责听懂和说明\n让确定性系统负责查询和计算", 0.65, 3.08, 3.1, 0.92, size=16, color="CDEBF2")
    add_pill(slide, "老板 / 甲方沟通版", 0.65, 5.73, 2.0, fill=BLUE, color=WHITE)
    add_text(slide, "2026-07-21", 0.65, 6.32, 2.0, 0.3, size=11, color="BFD5DE")
    # 右侧主视觉：业务问题进入受控闭环
    add_text(slide, "一句话价值", 5.05, 0.95, 2.0, 0.35, size=12, color=CYAN, bold=True)
    add_text(slide, "自然语言入口 + 可信计算内核 + 全链路证据", 5.05, 1.35, 7.3, 0.7, size=26, color=NAVY, bold=True)
    stages = [("问", "业务问题"), ("解", "结构化意图"), ("算", "受控工具"), ("证", "证据包"), ("核", "结果核验")]
    for i, (mark, label) in enumerate(stages):
        x = 5.05 + i * 1.48
        add_circle(slide, mark, x, 3.0, 0.72, fill=CYAN if i < 4 else GREEN)
        add_text(slide, label, x - 0.28, 3.85, 1.28, 0.42, size=11.5, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            add_line(slide, x + 0.76, 3.36, x + 1.36, 3.36, color=BORDER, width=2)
    add_rect(slide, 5.05, 5.0, 7.22, 0.92, fill=PALE, line=PALE)
    add_text(slide, "模型不能自由查库、编造数字或自动处置；生产动作保留人工责任门。", 5.35, 5.15, 6.65, 0.58, size=15, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    # 2 业务问题
    slide = prs.slides.add_slide(blank)
    add_header(slide, "传统“模型直连数据”路线，风险集中在事实与责任失控", "01｜业务问题", 2)
    add_card(slide, "业务需要", "自然语言提问\n跨指标组合分析\n快速解释与图表\n活动与异常辅助研判", 0.65, 1.55, 3.55, 4.8, accent=BLUE)
    add_card(slide, "直接接大模型的风险", "自由 SQL 或高成本查询\n指标口径混淆\n越权字段访问\n相关被误写成因果\n数字无法追溯\n运营建议责任不清", 4.55, 1.55, 3.75, 4.8, accent=RED)
    add_card(slide, "项目的回答", "模型只生成结构化候选\n工具白名单负责执行\n结果统一封装证据\nVerifier 做执行前/回答后硬门\n通知与处置由人确认", 8.65, 1.55, 4.05, 4.8, accent=GREEN)
    add_pill(slide, "目标不是“更自由的 Agent”，而是“更可信的业务协作系统”", 3.7, 6.53, 5.95, fill=NAVY, color=WHITE)
    add_footer(slide)

    # 3 方案
    slide = prs.slides.add_slide(blank)
    add_header(slide, "系统把语言、计算、证据与责任拆成可治理的分工", "02｜解决方案", 3)
    columns = [
        ("前端", "问与看", "业务问题\n回答 / 图表\n证据 / 轨迹", BLUE),
        ("状态机", "组织执行", "上下文\n意图 / 计划\n失败与回放", CYAN),
        ("工具内核", "真正计算", "QueryIR\n统计 / 预测\nGIS / 报告", ORANGE),
        ("质量门", "防止胡说", "EvidencePacket\nVerifier\n人工确认", GREEN),
    ]
    for i, (title, role, body, accent) in enumerate(columns):
        x = 0.65 + i * 3.12
        add_rect(slide, x, 1.65, 2.7, 4.2, fill=MIST, line=BORDER)
        add_circle(slide, str(i + 1), x + 0.22, 1.9, 0.48, fill=accent)
        add_text(slide, title, x + 0.82, 1.82, 1.55, 0.42, size=18, color=NAVY, bold=True)
        add_text(slide, role, x + 0.25, 2.58, 2.2, 0.42, size=13, color=accent, bold=True, align=PP_ALIGN.CENTER)
        add_line(slide, x + 0.35, 3.15, x + 2.35, 3.15, color=BORDER, width=1)
        add_text(slide, body, x + 0.3, 3.45, 2.1, 1.62, size=14, color=TEXT, align=PP_ALIGN.CENTER)
        if i < 3:
            add_line(slide, x + 2.72, 3.78, x + 3.0, 3.78, color=CYAN, width=2.5)
    add_rect(slide, 1.35, 6.18, 10.6, 0.64, fill=PALE, line=PALE)
    add_text(slide, "稳定契约：IntentEnvelope → TaskPlan → ToolResult → EvidencePacket / VerificationReport", 1.55, 6.28, 10.2, 0.42, size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)

    # 4 执行流程
    slide = prs.slides.add_slide(blank)
    add_header(slide, "一次回答经过 10 个可观察节点，失败可定位到具体步骤", "03｜运行闭环", 4)
    labels = [
        "接收问题", "构造上下文", "意图候选", "基线核验", "任务计划",
        "工具执行", "证据封装", "证据化回答", "结果核验", "轨迹返回",
    ]
    for i, label in enumerate(labels):
        row = i // 5
        col = i % 5
        x = 0.72 + col * 2.48
        y = 1.72 + row * 2.35
        accent = BLUE if i < 2 else CYAN if i < 5 else ORANGE if i < 7 else GREEN
        add_rect(slide, x, y, 2.05, 1.18, fill=WHITE, line=accent, line_width=1.5)
        add_circle(slide, str(i + 1), x + 0.12, y + 0.14, 0.4, fill=accent)
        add_text(slide, label, x + 0.57, y + 0.12, 1.33, 0.46, size=13.5, color=NAVY, bold=True)
        owner = "LLM 候选" if i in {2, 4, 7} else "确定性" if i not in {9} else "前端 / 审计"
        add_text(slide, owner, x + 0.25, y + 0.7, 1.55, 0.28, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
        if col < 4:
            add_line(slide, x + 2.05, y + 0.59, x + 2.37, y + 0.59, color=BORDER, width=2)
    add_text(slide, "工具失败时：只允许一次受限重规划；不重复成功步骤，不增加未知工具。", 1.9, 6.33, 9.55, 0.46, size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)

    # 5 能力地图
    slide = prs.slides.add_slide(blank)
    add_header(slide, "统一工具注册表，把查询、分析、预测和报告像积木一样组合", "04｜能力覆盖", 5)
    groups = [
        ("查询与比较", "指标目录｜QueryIR｜时段比较\n站点排序｜线路并行比较", BLUE),
        ("统计与诊断", "增长｜Pearson / 滞后相关\n异常｜趋势｜候选原因树", CYAN),
        ("预测与活动", "参考日基线｜活动规则\n时间序列｜回测样例｜SOP", ORANGE),
        ("跨网与空间", "轨道 / 公交换乘｜窗口匹配\n地理编码｜OD 热力｜通勤画像", GREEN),
        ("审计与报告", "EvidencePacket｜运行轨迹\n本地报告｜反馈与数据集门", NAVY),
        ("多端交付", "React Web｜FastAPI｜CLI\n微信小程序｜多端应用工程", BLUE),
    ]
    for i, (title, body, accent) in enumerate(groups):
        row, col = divmod(i, 3)
        add_card(slide, title, body, 0.7 + col * 4.2, 1.58 + row * 2.38, 3.78, 1.92, accent=accent)
    add_text(slide, "注意：跨网、GIS、活动、实时和 SOP 路径目前使用合成夹具。", 2.25, 6.46, 8.85, 0.42, size=13.5, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)

    # 6 场景
    slide = prs.slides.add_slide(blank)
    add_header(slide, "业务人员看到的不只是答案，还能看到答案如何产生", "05｜用户体验", 6)
    if SCREENSHOT.exists():
        slide.shapes.add_picture(str(SCREENSHOT), Inches(0.65), Inches(1.48), width=Inches(6.15), height=Inches(5.18))
    add_card(slide, "对老板 / 甲方", "看清系统当前能力、数据范围、模型是否真实调用、哪些仍待生产准入。", 7.2, 1.52, 5.42, 1.32, accent=BLUE)
    add_card(slide, "对运营人员", "自然语言提问；核验通过才展示为可采纳结果；高风险建议明确要求人工确认。", 7.2, 3.0, 5.42, 1.45, accent=GREEN)
    add_card(slide, "对开发人员", "一次 RunRecord 同时提供意图、计划、工具、证据、核验、模型 usage 与事件轨迹。", 7.2, 4.62, 5.42, 1.45, accent=CYAN)
    add_text(slide, "Web 智能分析页面（本地真实运行截图）", 0.85, 6.72, 5.75, 0.28, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide)

    # 7 治理
    slide = prs.slides.add_slide(blank)
    add_header(slide, "四道硬门把模型能力限制在可验证、可审计范围内", "06｜安全与治理", 7)
    gates = [
        ("意图门", "实体、指标、时间与任务类型\n必须匹配 protected intent", BLUE),
        ("计划门", "工具、参数、顺序与依赖\n必须匹配 protected plan", CYAN),
        ("证据门", "每个数字回到 ToolResult\n并获得 evidence_id", ORANGE),
        ("责任门", "通知、处置、生产数据\n保留人工最终确认", GREEN),
    ]
    for i, (title, body, accent) in enumerate(gates):
        x = 0.72 + i * 3.12
        add_rect(slide, x, 1.72, 2.72, 3.58, fill=MIST, line=accent, line_width=1.7)
        add_circle(slide, str(i + 1), x + 0.96, 2.03, 0.8, fill=accent)
        add_text(slide, title, x + 0.36, 3.03, 2.0, 0.48, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.3, 3.72, 2.12, 1.05, size=12.5, color=TEXT, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.05, 5.75, 11.2, 0.86, fill=NAVY, line=NAVY)
    add_text(slide, "系统当前不做：自由 SQL｜生产写库｜自动通知｜长期后台调度｜自动运营处置", 1.38, 5.9, 10.55, 0.55, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)

    # 8 验证
    slide = prs.slides.add_slide(blank)
    add_header(slide, "当前代码已形成可复跑质量基线，但不把工程验证冒充业务验收", "07｜真实验证", 8)
    metrics = [
        ("103", "Python 自动测试", GREEN),
        ("100/100", "智能体 Gold Cases", BLUE),
        ("11", "Web 浏览器 E2E", CYAN),
        ("0", "Ruff / Biome 新增错误", ORANGE),
    ]
    for i, (value, label, accent) in enumerate(metrics):
        x = 0.7 + i * 3.12
        add_rect(slide, x, 1.62, 2.75, 1.72, fill=WHITE, line=accent, line_width=1.5)
        add_text(slide, value, x + 0.2, 1.82, 2.35, 0.68, size=28, color=accent, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.2, 2.62, 2.35, 0.36, size=11.5, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    chart_data = ChartData()
    chart_data.categories = ["Python 测试", "Gold Cases", "Web E2E"]
    chart_data.add_series("通过数量", (103, 100, 11))
    chart = cast(Any, slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.72),
        Inches(3.7),
        Inches(4.45),
        Inches(2.22),
        chart_data,
    )).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 110
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.tick_labels.font.name = FONT
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.name = FONT
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = rgb(CYAN)
    table_shape = slide.shapes.add_table(3, 2, Inches(5.48), Inches(3.7), Inches(7.15), Inches(2.22))
    table = table_shape.table
    table.cell(0, 0).text = "工程验证已证明"
    table.cell(0, 1).text = "生产验收尚未证明"
    table.cell(1, 0).text = "契约、边界、工具组合、证据核验"
    table.cell(1, 1).text = "真实预测准确率、正式 SOP、运营责任"
    table.cell(2, 0).text = "前后端构建与主要交互路径可复跑"
    table.cell(2, 1).text = "权威数据权限、并发性能、生产模型端点"
    style_table(table)
    add_text(slide, "验证时间：2026-07-21｜结果来自当前本地工作树", 3.55, 6.3, 6.2, 0.36, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide)

    # 9 合作开发
    slide = prs.slides.add_slide(blank)
    add_header(slide, "稳定契约让甲方数据、算法与界面可以分工接入，而不重写全系统", "08｜合作开发", 9)
    lanes = [
        ("甲方业务与数据", "指标口径\n权威数据源\n权限与 SOP\n验收 Gold Cases", BLUE),
        ("算法与分析团队", "纯分析函数\n输入输出契约\n回测与误差\nToolRegistry 接入", ORANGE),
        ("平台开发团队", "FastAPI / OpenAPI\nWeb / 微信\n审计与部署\nCI / 安全", CYAN),
        ("联合治理", "EvidencePacket\nVerifier\n人工闸门\n生产准入", GREEN),
    ]
    for i, (title, body, accent) in enumerate(lanes):
        x = 0.7 + i * 3.12
        add_rect(slide, x, 1.58, 2.72, 4.55, fill=WHITE, line=BORDER)
        add_rect(slide, x, 1.58, 2.72, 0.68, fill=accent, line=accent, radius=False, line_width=0)
        add_text(slide, title, x + 0.18, 1.68, 2.36, 0.43, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.35, 2.75, 2.02, 2.35, size=14, color=TEXT, align=PP_ALIGN.CENTER)
        add_pill(slide, "明确 owner / verifier", x + 0.42, 5.45, 1.88, fill=PALE, color=accent)
    add_text(slide, "推荐第一步：共同确认指标字典、10–20 条真实脱敏 Gold Cases 和生产权限边界。", 1.7, 6.42, 9.95, 0.42, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)

    # 10 决策与路线
    slide = prs.slides.add_slide(blank)
    add_header(slide, "建议以“小范围只读试点”推进，用真实证据决定是否扩大模型与数据权限", "09｜下一步决策", 10)
    phases = [
        ("阶段 A", "联合定义", "指标字典\n数据责任人\nGold Cases\n准入/退出条件", BLUE),
        ("阶段 B", "只读试点", "脱敏权威数据\n最小 RBAC\n性能与审计\n人工采纳记录", CYAN),
        ("阶段 C", "评估净增益", "确定性基线对照\n错误归因\n预测准确率\n业务可用性", ORANGE),
        ("阶段 D", "受控扩展", "正式 SOP\n通知/调度审批\n灾备与回滚\n逐项开放权限", GREEN),
    ]
    for i, (phase, title, body, accent) in enumerate(phases):
        x = 0.65 + i * 3.14
        add_rect(slide, x, 1.52, 2.78, 4.55, fill=MIST, line=accent, line_width=1.5)
        add_pill(slide, phase, x + 0.62, 1.8, 1.55, fill=accent, color=WHITE)
        add_text(slide, title, x + 0.28, 2.45, 2.22, 0.5, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.36, 3.22, 2.06, 1.82, size=13, color=TEXT, align=PP_ALIGN.CENTER)
        if i < 3:
            add_line(slide, x + 2.8, 3.8, x + 3.04, 3.8, color=CYAN, width=2.5)
    add_rect(slide, 1.1, 6.28, 11.1, 0.54, fill=NAVY, line=NAVY)
    add_text(slide, "本次希望达成：确认试点场景、数据责任人、Gold Cases 与生产准入负责人。", 1.35, 6.35, 10.6, 0.38, size=14.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "决策建议｜先只读、先小样、先验证，再扩大权限")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    normalize_chart_axis_ids(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
